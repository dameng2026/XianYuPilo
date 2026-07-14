import datetime
from io import BytesIO
from typing import Any, Optional

from fastapi import APIRouter, Body, Depends, File, HTTPException, Query, UploadFile
from sqlalchemy import func, or_, select
from sqlalchemy.ext.asyncio import AsyncSession

from ....core.database import get_db
from ....core.response import ResultObject
from ....models.entities import AutoReplyRule, QuickReplyTemplate
from ....services.ai_provider import _resolve_ai_config, generate_text, is_ai_configured
from ....services.business_settings import (
    AI_CS_SETTING_KEY,
    ALLOWED_BUSINESS_SETTING_CATEGORIES,
    BusinessSettingValidationError,
    build_default_business_setting,
    load_business_setting,
    save_business_setting,
)
from ..deps import get_current_user

router = APIRouter(tags=["businessSettingsCompat"])
AI_SETTINGS_HINT = "未配置通用模型，请先前往系统设置中的“模型配置”填写 baseUrl、apiKey 与模型名称。"


def _format_datetime(value: Optional[datetime.datetime]) -> Optional[str]:
    if value is None:
        return None
    return value.isoformat(sep=" ", timespec="seconds")


def _rule_to_record(rule: AutoReplyRule) -> dict[str, Any]:
    return {
        "id": rule.id,
        "accountId": rule.account_id,
        "xianyuAccountId": rule.account_id,
        "ruleName": rule.rule_name or "",
        "matchType": rule.match_type or "keyword",
        "matchKeywords": rule.match_keywords or "",
        "replyContent": rule.reply_content or "",
        "replyMode": rule.reply_mode or rule.match_type or "keyword",
        "status": rule.status if rule.status is not None else 1,
        "priority": rule.priority or 0,
        "createdTime": _format_datetime(rule.created_time),
        "updatedTime": _format_datetime(rule.updated_time),
    }


def _template_to_record(template: QuickReplyTemplate) -> dict[str, Any]:
    return {
        "id": template.id,
        "accountId": template.account_id,
        "title": template.title,
        "content": template.content,
        "text": template.content,
        "sortOrder": template.sort_order or 0,
        "status": template.status if template.status is not None else 1,
        "createdAt": _format_datetime(template.created_time),
        "updatedAt": _format_datetime(template.updated_time),
    }


def _pick_keywords(raw: Any) -> list[str]:
    text_value = str(raw or "").strip()
    if not text_value:
        return []
    separators = [",", "，", "\n", "\r", ";", "；", "|", "/", "、"]
    for separator in separators[1:]:
        text_value = text_value.replace(separator, separators[0])
    return [item.strip() for item in text_value.split(separators[0]) if item.strip()]


async def _decode_text_upload(file_name: str, content: bytes) -> str:
    suffix = file_name.lower().rsplit(".", 1)[-1] if "." in file_name else ""
    if suffix in {"md", "txt", "csv"}:
        for encoding in ("utf-8-sig", "utf-8", "gb18030"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="ignore")

    if suffix == "pptx":
        try:
            from pptx import Presentation

            presentation = Presentation(BytesIO(content))
            parts: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    text_value = getattr(shape, "text", "")
                    if text_value:
                        parts.append(str(text_value).strip())
            return "\n".join([part for part in parts if part])
        except Exception:  # noqa: BLE001
            return ""

    if suffix == "xlsx":
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
            parts: list[str] = []
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value not in (None, "")]
                    if values:
                        parts.append(" | ".join(values))
            return "\n".join(parts)
        except Exception:  # noqa: BLE001
            return ""

    return ""


@router.get("/business-settings/ai-customer-service/defaults", response_model=ResultObject)
async def get_ai_customer_service_defaults(
    _: dict = Depends(get_current_user),
):
    return ResultObject.success(build_default_business_setting(AI_CS_SETTING_KEY))


@router.post("/business-settings/ai-customer-service/test", response_model=ResultObject)
async def test_ai_customer_service(
    body: dict[str, Any] = Body(default_factory=dict),
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    config = await load_business_setting(db, AI_CS_SETTING_KEY)
    system_prompt = str(config.get("systemPrompt") or "").strip()[:6_000]
    raw_message = body.get("message")
    user_message = (
        "你好，这个商品还可以再优惠一点吗？"
        if raw_message is None
        else str(raw_message).strip()
    )
    if not user_message:
        raise HTTPException(status_code=422, detail="测试消息不能为空")
    if len(user_message) > 4_000:
        raise HTTPException(status_code=422, detail="测试消息不能超过 4000 个字符")

    provider_config = await _resolve_ai_config()
    provider_configured = is_ai_configured(provider_config)
    if not provider_configured:
        raise HTTPException(status_code=503, detail=AI_SETTINGS_HINT)

    try:
        result = await generate_text(
            scene="ai_customer_service_test",
            system_prompt=system_prompt,
            user_prompt=user_message,
            temperature=0.6,
        )
    except Exception as exc:  # noqa: BLE001
        # Provider errors can contain URLs, request bodies, or credentials.
        # The API exposes only a stable unavailable state.
        del exc
        raise HTTPException(
            status_code=503,
            detail="AI 模型调用暂不可用，请检查模型配置与网络后重试",
        ) from None
    if isinstance(result, dict) and result.get("ok") and result.get("content"):
        return ResultObject.success({
            "ok": True,
            "reply": str(result.get("content") or "").strip(),
            "configured": True,
            "usage": result.get("usage") or {},
        })

    raise HTTPException(
        status_code=503,
        detail="AI 模型未返回有效回复，请检查模型服务状态后重试",
    )


@router.post("/business-settings/ai-customer-service/upload-knowledge", response_model=ResultObject)
async def upload_ai_customer_service_knowledge(
    file: UploadFile = File(...),
    _: dict = Depends(get_current_user),
):
    file_name = file.filename or ""
    if not file_name.strip():
        return ResultObject.failed("文件名不能为空", code=400)

    suffix = f".{file_name.lower().rsplit('.', 1)[-1]}" if "." in file_name else ""
    legacy_unsupported = {".ppt", ".xls"}
    if suffix in legacy_unsupported:
        raise HTTPException(
            status_code=415,
            detail="不支持旧版二进制 Office 格式，请另存为 .pptx / .xlsx 后重新上传",
        )

    allowed = {".md", ".txt", ".pptx", ".xlsx", ".csv"}
    if suffix not in allowed:
        return ResultObject.failed(f"不支持的文件格式：{suffix}", code=400)

    content = await file.read()
    if not content:
        return ResultObject.failed("上传文件不能为空", code=400)

    if len(content) > 10 * 1024 * 1024:
        return ResultObject.failed("文件不能超过 10MB", code=400)

    extracted_text = await _decode_text_upload(file_name, content)
    if not extracted_text.strip():
        return ResultObject.failed("未能从文件中提取有效内容", code=400)

    lines = [line.strip() for line in extracted_text.splitlines() if line.strip()]
    return ResultObject.success({
        "fileName": file_name,
        "extractedText": extracted_text[:20000],
        "ruleCount": len(lines),
    })


@router.get("/business-settings/{category}", response_model=ResultObject)
async def get_business_settings(
    category: str,
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    if category not in ALLOWED_BUSINESS_SETTING_CATEGORIES:
        return ResultObject.failed(f"不支持的配置分类: {category}", code=400)
    return ResultObject.success(await load_business_setting(db, category))


@router.post("/business-settings/{category}", response_model=ResultObject)
async def save_business_settings(
    category: str,
    body: dict[str, Any] = Body(default_factory=dict),
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    if category not in ALLOWED_BUSINESS_SETTING_CATEGORIES:
        return ResultObject.failed(f"不支持的配置分类: {category}", code=400)
    try:
        saved = await save_business_setting(db, category, dict(body or {}))
    except BusinessSettingValidationError as exc:
        raise HTTPException(status_code=422, detail=str(exc)) from None
    return ResultObject.success(saved)


@router.get("/quick-reply/templates", response_model=ResultObject)
async def list_quick_reply_templates(
    account_id: Optional[int] = Query(None, alias="accountId"),
    size: int = Query(100, ge=1, le=500),
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    query = select(QuickReplyTemplate).where(QuickReplyTemplate.deleted == 0)
    if account_id is None:
        query = query.where(
            or_(QuickReplyTemplate.account_id.is_(None), QuickReplyTemplate.account_id == 0)
        )
    else:
        query = query.where(
            or_(QuickReplyTemplate.account_id.is_(None), QuickReplyTemplate.account_id == account_id)
        )

    result = await db.execute(
        query.order_by(QuickReplyTemplate.sort_order.asc(), QuickReplyTemplate.id.asc()).limit(size)
    )
    items = [_template_to_record(item) for item in result.scalars().all()]
    return ResultObject.success({"records": items, "total": len(items)})


@router.post("/quick-reply/templates", response_model=ResultObject)
async def save_quick_reply_template(
    body: dict[str, Any] = Body(default_factory=dict),
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    title = str(body.get("title") or "").strip()
    content = str(body.get("content") or "").strip()
    if not title or not content:
        return ResultObject.failed("模板标题和内容不能为空", code=400)

    template_id = body.get("id")
    account_id = body.get("accountId")
    sort_order = int(body.get("sortOrder") or body.get("sort_order") or 0)
    if template_id:
        result = await db.execute(
            select(QuickReplyTemplate).where(
                QuickReplyTemplate.id == int(template_id),
                QuickReplyTemplate.deleted == 0,
            )
        )
        template = result.scalar_one_or_none()
        if not template:
            return ResultObject.failed("模板不存在", code=404)
        template.title = title
        template.content = content
        template.sort_order = sort_order
        if account_id is not None:
            template.account_id = int(account_id)
        await db.commit()
        await db.refresh(template)
        return ResultObject.success(_template_to_record(template), "模板已更新")

    template = QuickReplyTemplate(
        account_id=int(account_id) if account_id is not None else None,
        title=title,
        content=content,
        sort_order=sort_order,
        status=1,
        deleted=0,
    )
    db.add(template)
    await db.commit()
    await db.refresh(template)
    return ResultObject.success(_template_to_record(template), "模板已创建")


@router.delete("/quick-reply/templates/{template_id}", response_model=ResultObject)
async def delete_quick_reply_template(
    template_id: int,
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    result = await db.execute(
        select(QuickReplyTemplate).where(
            QuickReplyTemplate.id == template_id,
            QuickReplyTemplate.deleted == 0,
        )
    )
    template = result.scalar_one_or_none()
    if not template:
        return ResultObject.failed("模板不存在", code=404)
    template.deleted = 1
    await db.commit()
    return ResultObject.success({"id": template_id}, "模板已删除")


@router.get("/auto-reply/rules", response_model=ResultObject)
async def list_auto_reply_rules(
    account_id: Optional[int] = Query(None, alias="accountId"),
    current: int = Query(1, ge=1),
    size: int = Query(10, ge=1, le=500),
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    query = select(AutoReplyRule).where(AutoReplyRule.deleted == 0)
    if account_id is not None:
        query = query.where(AutoReplyRule.account_id == account_id)

    total = (await db.execute(select(func.count()).select_from(query.subquery()))).scalar() or 0
    result = await db.execute(
        query.order_by(AutoReplyRule.priority.desc(), AutoReplyRule.id.desc())
        .offset((current - 1) * size)
        .limit(size)
    )
    records = [_rule_to_record(item) for item in result.scalars().all()]
    return ResultObject.success({
        "records": records,
        "total": total,
        "current": current,
        "size": size,
    })


@router.post("/auto-reply/rules", response_model=ResultObject)
async def create_auto_reply_rule(
    body: dict[str, Any] = Body(default_factory=dict),
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    rule = AutoReplyRule(
        account_id=body.get("accountId") or body.get("xianyuAccountId"),
        rule_name=str(body.get("ruleName") or body.get("rule_name") or "").strip() or "自动回复规则",
        match_type=str(body.get("matchType") or body.get("match_type") or "keyword").strip() or "keyword",
        match_keywords=str(body.get("matchKeywords") or body.get("match_keywords") or "").strip() or None,
        reply_content=str(body.get("replyContent") or body.get("reply_content") or "").strip() or None,
        reply_mode=str(body.get("replyMode") or body.get("reply_mode") or body.get("matchType") or "keyword").strip() or "keyword",
        status=int(body.get("status") if body.get("status") is not None else 1),
        priority=int(body.get("priority") or 0),
        deleted=0,
    )
    db.add(rule)
    await db.commit()
    await db.refresh(rule)
    return ResultObject.success(_rule_to_record(rule), "规则已创建")


@router.put("/auto-reply/rules/{rule_id}", response_model=ResultObject)
async def update_auto_reply_rule(
    rule_id: int,
    body: dict[str, Any] = Body(default_factory=dict),
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    result = await db.execute(
        select(AutoReplyRule).where(AutoReplyRule.id == rule_id, AutoReplyRule.deleted == 0)
    )
    rule = result.scalar_one_or_none()
    if not rule:
        return ResultObject.failed("规则不存在", code=404)

    if "accountId" in body or "xianyuAccountId" in body:
        rule.account_id = body.get("accountId") or body.get("xianyuAccountId")
    if "ruleName" in body or "rule_name" in body:
        rule.rule_name = str(body.get("ruleName") or body.get("rule_name") or "").strip() or rule.rule_name
    if "matchType" in body or "match_type" in body:
        rule.match_type = str(body.get("matchType") or body.get("match_type") or "keyword").strip() or "keyword"
    if "matchKeywords" in body or "match_keywords" in body:
        rule.match_keywords = str(body.get("matchKeywords") or body.get("match_keywords") or "").strip() or None
    if "replyContent" in body or "reply_content" in body:
        rule.reply_content = str(body.get("replyContent") or body.get("reply_content") or "").strip() or None
    if "replyMode" in body or "reply_mode" in body:
        rule.reply_mode = str(body.get("replyMode") or body.get("reply_mode") or rule.reply_mode or "keyword").strip() or "keyword"
    if "status" in body:
        rule.status = int(body.get("status") if body.get("status") is not None else 1)
    if "priority" in body:
        rule.priority = int(body.get("priority") or 0)

    await db.commit()
    await db.refresh(rule)
    return ResultObject.success(_rule_to_record(rule), "规则已更新")


@router.delete("/auto-reply/rules/{rule_id}", response_model=ResultObject)
async def delete_auto_reply_rule(
    rule_id: int,
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    result = await db.execute(
        select(AutoReplyRule).where(AutoReplyRule.id == rule_id, AutoReplyRule.deleted == 0)
    )
    rule = result.scalar_one_or_none()
    if not rule:
        return ResultObject.failed("规则不存在", code=404)
    rule.deleted = 1
    await db.commit()
    return ResultObject.success({"id": rule_id}, "规则已删除")


@router.post("/auto-reply/rules/preview", response_model=ResultObject)
async def preview_auto_reply_rule(
    body: dict[str, Any] = Body(default_factory=dict),
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    account_id = body.get("accountId") or body.get("xianyuAccountId")
    message = str(body.get("message") or "").strip()
    query = select(AutoReplyRule).where(
        AutoReplyRule.deleted == 0,
        AutoReplyRule.status == 1,
    )
    if account_id is not None:
        query = query.where(or_(AutoReplyRule.account_id == int(account_id), AutoReplyRule.account_id.is_(None)))

    result = await db.execute(query.order_by(AutoReplyRule.priority.desc(), AutoReplyRule.id.desc()))
    rules = result.scalars().all()
    matched_rule = None
    for rule in rules:
        if (rule.match_type or "keyword") == "all":
            matched_rule = rule
            break
        keywords = _pick_keywords(rule.match_keywords)
        if keywords and any(keyword.lower() in message.lower() for keyword in keywords):
            matched_rule = rule
            break

    if not matched_rule:
        return ResultObject.success({
            "matched": False,
            "replyContent": "",
            "rule": None,
        })

    return ResultObject.success({
        "matched": True,
        "replyContent": matched_rule.reply_content or "",
        "rule": _rule_to_record(matched_rule),
    })


@router.get("/auto-reply/rules/logs", response_model=ResultObject)
async def get_auto_reply_logs(
    current: int = Query(1, ge=1),
    size: int = Query(20, ge=1, le=500),
    _: dict = Depends(get_current_user),
):
    return ResultObject.success({
        "records": [],
        "total": 0,
        "current": current,
        "size": size,
    })


@router.get("/auto-reply/rules/stats", response_model=ResultObject)
async def get_auto_reply_stats(
    days: int = Query(7, ge=1, le=365),
    db: AsyncSession = Depends(get_db),
    _: dict = Depends(get_current_user),
):
    total = (
        await db.execute(
            select(func.count()).select_from(AutoReplyRule).where(AutoReplyRule.deleted == 0)
        )
    ).scalar() or 0
    enabled = (
        await db.execute(
            select(func.count()).select_from(AutoReplyRule).where(
                AutoReplyRule.deleted == 0,
                AutoReplyRule.status == 1,
            )
        )
    ).scalar() or 0
    return ResultObject.success({
        "days": days,
        "totalRules": total,
        "enabledRules": enabled,
        "disabledRules": max(total - enabled, 0),
        "matchedCount": 0,
        "replyCount": 0,
    })
